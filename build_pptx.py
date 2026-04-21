"""Populate PTP_LAN8651.pptx — keeps Microchip corporate template (logos, master)
and fills it with our PTP / 10BASE-T1S / LAN8651 / results content."""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PATH = r"C:\work\ptp\check\net_10base_t1s\PTP_LAN8651.pptx"

prs = Presentation(PATH)

# Remove the 4 sample slides that ship with the template
sldIdLst = prs.slides._sldIdLst
part = prs.part
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    part.drop_rel(rId)
    sldIdLst.remove(sldId)

SW, SH = prs.slide_width, prs.slide_height  # 12188825 x 6858000 (~13.33 x 7.5 in)

LAYOUT = {l.name: l for l in prs.slide_layouts}
L_TITLE   = LAYOUT["Title Slide"]
L_CONTENT = LAYOUT["Title and Content"]
L_TITLE_ONLY = LAYOUT["Title Only"]
L_SECTION = LAYOUT["Section Header"]

# accent colours — kept moderate so the Microchip template drives the overall look
NAVY = RGBColor(0x0B, 0x2E, 0x4F)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xE8, 0xEE, 0xF5)
OK = RGBColor(0x2C, 0xA0, 0x2C)
WARN = RGBColor(0xD6, 0x27, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def get_placeholder(slide, name_substr):
    for ph in slide.placeholders:
        if name_substr.lower() in ph.name.lower():
            return ph
    return None


def set_title(slide, text):
    ph = get_placeholder(slide, "Title")
    if ph is None:
        return None
    ph.text_frame.text = text
    return ph


def fill_content(slide, items, size=18):
    """Fill the Content placeholder with a bullet list. items: list of str or (level, str)."""
    ph = get_placeholder(slide, "Content")
    if ph is None:
        ph = get_placeholder(slide, "Text Placeholder")
    if ph is None:
        return None
    tf = ph.text_frame
    tf.clear()
    first = True
    for item in items:
        level, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size - 2*level)
    return ph


def add_slide(layout, title=None, content=None, content_size=18):
    s = prs.slides.add_slide(layout)
    if title is not None:
        set_title(s, title)
    if content is not None:
        fill_content(s, content, size=content_size)
    return s


def box(slide, l, t, w, h, text, fill=ACCENT, fc=WHITE, size=14, bold=True, italic=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = 1
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = fc
    return sh


def line(slide, x1, y1, x2, y2, color=GREY, width=2, arrow=False):
    shape_id = 2 if arrow else 1
    c = slide.shapes.add_connector(shape_id, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def label(slide, x, y, w, h, text, size=12, color=NAVY, bold=False, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


# ============================================================================
# 1. Title slide
# ============================================================================
s = prs.slides.add_slide(L_TITLE)
ph_title = get_placeholder(s, "Title")
if ph_title is not None:
    ph_title.text_frame.text = "PTP over 10BASE-T1S"
# Subtitle / text placeholders
for name, txt in [
    ("Text Placeholder 2", "Theory · PHY Context · LAN8651 Specifics · Measurement Results"),
    ("Text Placeholder 3", "Martin Ruppert  |  Microchip WNET  |  2026-04-21"),
]:
    ph = get_placeholder(s, name)
    if ph is not None:
        ph.text_frame.text = txt

# ============================================================================
# 2. Overview for a lay audience — What the demo shows
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "What the Demo Shows")

# Two boards with a single wire, a shared clock above, aligned pulses below
box(s, 1.3, 1.8, 3.6, 1.3, "Board A\nATSAME54 + LAN8651", fill=NAVY, size=16)
box(s, 8.5, 1.8, 3.6, 1.3, "Board B\nATSAME54 + LAN8651", fill=NAVY, size=16)
box(s, 4.95, 2.15, 3.4, 0.6, "single-pair Ethernet  (10BASE-T1S)", fill=OK, size=13)
# link lines
line(s, 4.9, 2.45, 5.0, 2.45, color=OK, width=3)
line(s, 8.35, 2.45, 8.5, 2.45, color=OK, width=3)

# "Shared clock" banner
box(s, 3.5, 3.3, 6.3, 0.7, "↕  Shared notion of time  (PTP, IEEE 1588)  ↕", fill=ACCENT, size=16)

# Left side: without sync (drifting pulses) — just a caption + small schematic idea
label(s, 0.6, 4.25, 6.2, 0.4, "Without synchronisation", size=15, bold=True, color=WARN)
# drifting pulses (small rectangles on a line, offset between two rows)
for x in [0.8, 2.1, 3.4, 4.7]:
    box(s, x,     4.75, 0.5, 0.35, "", fill=NAVY, size=1)
for x in [1.05, 2.45, 3.9, 5.4]:
    box(s, x,     5.25, 0.5, 0.35, "", fill=NAVY, size=1)
label(s, 0.6, 5.7, 6.2, 0.35, "Pulses on the two boards slowly drift apart.",
      size=12, italic=True, color=GREY)

# Right side: with sync (aligned pulses)
label(s, 7.0, 4.25, 6.2, 0.4, "With PTP synchronisation", size=15, bold=True, color=OK)
for x in [7.2, 8.5, 9.8, 11.1]:
    box(s, x,     4.75, 0.5, 0.35, "", fill=OK, size=1)
for x in [7.2, 8.5, 9.8, 11.1]:
    box(s, x,     5.25, 0.5, 0.35, "", fill=OK, size=1)
label(s, 7.0, 5.7, 6.2, 0.35, "Pulses on both boards line up on the wire — within a few microseconds.",
      size=12, italic=True, color=GREY)

# Bottom takeaway banner
box(s, 0.6, 6.3, 12.1, 0.6,
    "Goal: two independent microcontrollers that agree on time well enough to act together.",
    fill=LIGHT, fc=NAVY, size=15, italic=True)

# ============================================================================
# 3. Overview for a lay audience — How we achieved it
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "How We Achieved It")

# Four-step pipeline, left to right
step_w, step_h = 2.9, 1.7
step_y = 1.9
gap = 0.15
xs = [0.35, 0.35 + (step_w + gap),
      0.35 + 2*(step_w + gap), 0.35 + 3*(step_w + gap)]

steps = [
    ("1.  Hardware\ntimestamping",
     "LAN8651 latches the exact moment a packet crosses the wire."),
    ("2.  PTP protocol\n(IEEE 1588)",
     "Boards exchange timestamps and compute offset + path delay."),
    ("3.  Software\nwall-clock",
     "A shared nanosecond clock runs on every board."),
    ("4.  Cyclic Fire\ndemo",
     "Each board toggles a GPIO at scheduled clock values."),
]
for i, (title, _) in enumerate(steps):
    box(s, xs[i], step_y, step_w, step_h, title,
        fill=NAVY if i % 2 == 0 else ACCENT, size=15)

# Arrows between the four boxes
for i in range(3):
    x1 = xs[i] + step_w
    x2 = xs[i+1]
    line(s, x1, step_y + step_h/2, x2, step_y + step_h/2,
         color=GREY, width=3, arrow=True)

# Explanation under each step
for i, (_, expl) in enumerate(steps):
    tb = s.shapes.add_textbox(Inches(xs[i]), Inches(step_y + step_h + 0.1),
                              Inches(step_w), Inches(1.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = expl
    r.font.size = Pt(12); r.font.color.rgb = GREY

# Outcome banner
box(s, 0.35, 6.1, 12.25, 0.8,
    "Result: a cross-board edge delta of ~10 µs median — inside the 100 µs product window.",
    fill=OK, size=16)

# ============================================================================
# 4. Standalone Button-LED Demo — functional description
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Standalone Button-LED Demo — How It Works")

# Intro strip
box(s, 0.4, 1.45, 12.5, 0.5,
    "Self-contained: two boards, a T1S cable, two buttons (SW1 / SW2), two LEDs (LED1 / LED2) — no PC, no CLI.",
    fill=LIGHT, fc=NAVY, size=13, italic=True)

# Three columns: Phase 1 / Phase 2 / Phase 3
col_w = 4.05
col_h = 3.7
col_y = 2.15
col_xs = [0.4, 4.65, 8.90]

# Phase headers
phase_titles = [
    "1.  Boot  —  Unsynchronised",
    "2.  Role Selection via Buttons",
    "3.  Synchronised  —  Lock-Step",
]
phase_colors = [WARN, ACCENT, OK]
for i, (t, col) in enumerate(zip(phase_titles, phase_colors)):
    box(s, col_xs[i], col_y, col_w, 0.55, t, fill=col, size=14)

# Phase 1 body
p1 = [
    "• Both boards boot with PTP OFF.",
    "• LED1 blinks at 1 Hz on each board (500 ms on / 500 ms off).",
    "• LED2 is off.",
    "• Because the two crystals differ by ~100 ppm, the two LED1 blinks visibly drift apart.",
    "• PD10 mirrors LED1 — a scope can watch the drift quantitatively.",
]
tb = s.shapes.add_textbox(Inches(col_xs[0]+0.05), Inches(col_y+0.65),
                          Inches(col_w-0.1), Inches(col_h-0.7))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(p1):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(12); r.font.color.rgb = NAVY

# Phase 2 body
p2 = [
    "• One board: press SW1 → becomes PTP Follower.",
    "• Other board: press SW2 → becomes PTP Master.",
    "• On each board, LED2 starts blinking fast (4 Hz) to signal 'syncing'.",
    "• The Follower's servo runs until it reaches FINE lock.",
    "• The Master waits ~2 s (visual symmetry with the Follower).",
]
tb = s.shapes.add_textbox(Inches(col_xs[1]+0.05), Inches(col_y+0.65),
                          Inches(col_w-0.1), Inches(col_h-0.7))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(p2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(12); r.font.color.rgb = NAVY

# Phase 3 body
p3 = [
    "• Master:   LED2 solid ON  (full brightness).",
    "• Follower: LED2 noticeably dimmer  (~6 % PWM @ 250 Hz, no flicker).",
    "• → The darker LED2 is the 'this is the follower' marker — visible from across the room.",
    "• LED1 on both boards now blinks in lock-step at 1 Hz — drift is gone.",
    "• PD10 on a scope shows the two 1 Hz edges co-aligned within microseconds.",
]
tb = s.shapes.add_textbox(Inches(col_xs[2]+0.05), Inches(col_y+0.65),
                          Inches(col_w-0.1), Inches(col_h-0.7))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(p3):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(12); r.font.color.rgb = NAVY

# Bottom pin/reference strip
box(s, 0.4, 6.15, 12.5, 0.45,
    "SW1 = PD00   SW2 = PD01   LED1 = PC21 (1 Hz blink)   LED2 = PA16 (status)   "
    "PD10 = LED1 mirror (scope)",
    fill=GREY, fc=WHITE, size=12)
label(s, 0.4, 6.65, 12.5, 0.35,
      "Source: apps/tcpip_iperf_lan865x/firmware/src/standalone_demo.c  (branch: ptp-standlone-demo)",
      size=11, italic=True, color=GREY)

# ============================================================================
# 5. Standalone Demo — Why the LEDs Appear Synchronous
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Standalone Demo — Why the LEDs Blink in Sync")

# Chain diagram across the top:
#   PTP_CLOCK  →  cyclic_fire  →  decimator /2000  →  LED1
box(s, 0.3, 1.55, 2.9, 1.1,
    "PTP software\nwall-clock", fill=NAVY, size=14)
box(s, 3.6, 1.55, 2.9, 1.1,
    "cyclic_fire\ncallback\nevery 250 µs", fill=ACCENT, size=14)
box(s, 6.9, 1.55, 2.9, 1.1,
    "Decimator\n÷ 2000", fill=ACCENT, size=14)
box(s, 10.2, 1.55, 2.8, 1.1,
    "LED1 toggles\nevery 500 ms\n(= 1 Hz blink)", fill=OK, size=14)
for x in [3.2, 6.5, 9.8]:
    line(s, x, 2.1, x+0.4, 2.1, color=GREY, width=3, arrow=True)

# Small caption for the rates
label(s, 0.3, 2.7, 13.0, 0.35,
      "4 kHz callback  →  ÷ 2000  →  1 Hz visible blink.  "
      "Every step is derived from the same PTP wall-clock.",
      size=13, italic=True, color=GREY)

# Two columns: without PTP  vs  with PTP
col_y = 3.25
col_h = 2.8
box(s, 0.3, col_y, 6.25, 0.55, "Without PTP lock  (drifting)", fill=WARN, size=14)
box(s, 6.75, col_y, 6.25, 0.55, "With PTP lock  (synchronised)", fill=OK, size=14)

left = [
    "• Each board's callback is driven only by its own TC0 crystal.",
    "• Two crystals differ by ~100 ppm  →  the 4 kHz callback instants drift.",
    "• The ÷ 2000 divider faithfully inherits that drift.",
    "• LED1 on the two boards therefore blinks a little slower / faster.",
    "• Over 10 s the LED1 edges separate by tens of ms — clearly visible.",
]
right = [
    "• Both boards' PTP_CLOCK anchors are pinned to the same wall-clock.",
    "• cyclic_fire schedules each callback at an absolute PTP-ns time.",
    "• The 4 kHz callback fires at the identical PTP instant on both boards.",
    "• The ÷ 2000 divider converts that into a common 1 Hz LED edge.",
    "• LED1 on both boards toggles together — the human eye sees one blink.",
]
tb = s.shapes.add_textbox(Inches(0.35), Inches(col_y+0.65),
                          Inches(6.2), Inches(col_h-0.7))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(left):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(13); r.font.color.rgb = NAVY

tb = s.shapes.add_textbox(Inches(6.80), Inches(col_y+0.65),
                          Inches(6.2), Inches(col_h-0.7))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(right):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(13); r.font.color.rgb = NAVY

# Bottom takeaway banner
box(s, 0.3, 6.3, 12.7, 0.65,
    "Key idea: decimating a fast, PTP-anchored callback is what turns sub-µs network time "
    "into a slow, human-visible LED blink.",
    fill=LIGHT, fc=NAVY, size=14, italic=True)

# ============================================================================
# 6. Standalone Demo — Live Flow, Role Indicator & Auto-Recovery
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Standalone Demo — Live Flow & Auto-Recovery")

# Row 1: role indicator via LED2 brightness
box(s, 0.3, 1.55, 12.7, 0.55,
    "Role indicator  —  LED2 brightness tells master from follower at a glance",
    fill=NAVY, size=14)
label(s, 0.5, 2.15, 6.0, 0.4, "Master board", size=14, bold=True, color=OK)
label(s, 0.5, 2.5, 6.2, 0.4,  "LED2 = solid ON  (full brightness)",
      size=13, color=NAVY)
label(s, 6.8, 2.15, 6.0, 0.4, "Follower board", size=14, bold=True, color=ACCENT)
label(s, 6.8, 2.5, 6.2, 0.4,
      "LED2 = ~6 % PWM @ 250 Hz  (dimmed, no flicker)",
      size=13, color=NAVY)

# Row 2: SW2 master on/off toggle
box(s, 0.3, 3.05, 12.7, 0.55,
    "SW2 is a master on/off toggle  —  lets you live-demonstrate sync-loss and recovery",
    fill=NAVY, size=14)
sw2_lines = [
    "• Press SW2 on the master again → master stops sending Sync, returns to FREE.",
    "• Follower notices no Sync for > 1000 ms → enters DEMO_LOST → LED2 turns OFF.",
    "• Press SW2 on the master once more → master re-arms → follower auto-recovers → LED2 dimmed again.",
    "• SW1 stays one-shot: the follower self-recovers, no manual off-switch needed.",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(3.65), Inches(12.3), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(sw2_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(13); r.font.color.rgb = NAVY

# Row 3: recovery mechanisms
box(s, 0.3, 5.1, 12.7, 0.55,
    "Three recovery mechanisms cover the failure modes seen on the bench",
    fill=ACCENT, size=14)
rec_lines = [
    "1.  Stale sync-sequence-id  →  PTP_FOL_Reset() on DEMO_LOST entry so the first fresh Sync is accepted.",
    "2.  PTP_CLOCK backward-jump at restart  →  cyclic_fire watchdog samples cycle count every 500 ms "
    "and re-arms with a fresh 'now+period' anchor if the chain stalls.",
    "3.  Grandmaster liveness  →  PTP_FOL_GetLastSyncTick() latched at the top of processSync() — "
    "even rejected Syncs count as 'heard from GM'.",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
for i, l in enumerate(rec_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = l
    r.font.size = Pt(12); r.font.color.rgb = NAVY

label(s, 0.3, 7.15, 12.7, 0.35,
      "Companion: standalone_demo_test.py walks the operator through the run + gates the result with a Saleae capture.",
      size=11, italic=True, color=GREY)

# ============================================================================
# 7. Agenda
# ============================================================================
add_slide(L_CONTENT, "Agenda", [
    "Part 0 — Overview & standalone button-LED demo",
    "Part 1 — PTP theory (IEEE 1588)",
    (1, "Clock model, message exchange, timestamp accuracy, error sources"),
    "Part 2 — 10BASE-T1S context",
    (1, "Single-pair multidrop PHY, PLCA, implications for PTP"),
    "Part 3 — LAN8651 specifics",
    (1, "SPI MAC-PHY, timestamp engine, SFD-to-IRQ latency"),
    "Part 4 — Cyclic Fire application",
    (1, "Concept, use cases, API and patterns"),
    "Part 5 — PTP_CLOCK for firmware event logging",
    (1, "API, usage pattern, cross-board correlation"),
    "Part 6 — Measurement results, open issues & outlook",
    "Part 7 — Reproducing the project on your own machine",
])

# ============================================================================
# 3. PTP Theory — Purpose and clock model
# ============================================================================
add_slide(L_CONTENT, "PTP — Purpose & Clock Model", [
    "Goal: distribute a common notion of time across networked nodes",
    (1, "Target accuracy: sub-microsecond on a LAN — not reachable with NTP"),
    "Roles within a PTP domain:",
    (1, "Grandmaster (GM) — reference clock, source of truth"),
    (1, "Ordinary / Boundary clock — follower, steers local clock to GM"),
    (1, "Transparent clock — forwards and corrects for residence time"),
    "Best Master Clock Algorithm (BMCA) elects the GM automatically",
    "Local clock corrected in two dimensions:",
    (1, "Phase — absolute offset to GM (ns)"),
    (1, "Rate — frequency ratio between local and GM oscillator (ppb)"),
])

# ============================================================================
# 4. PTP Theory — Message exchange with diagram
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "PTP — Message Exchange")

gm_x, fol_x = 2.3, 10.5
top_y, bot_y = 2.0, 5.9

line(s, gm_x, top_y, gm_x, bot_y, color=NAVY, width=3)
line(s, fol_x, top_y, fol_x, bot_y, color=NAVY, width=3)

box(s, gm_x-1.1, 1.55, 2.2, 0.4, "Grandmaster", fill=NAVY, size=14)
box(s, fol_x-1.1, 1.55, 2.2, 0.4, "Follower", fill=NAVY, size=14)


def msg(x1, y1, x2, y2, lbl, tx, ty, color=ACCENT):
    line(s, x1, y1, x2, y2, color=color, width=2, arrow=True)
    label(s, tx, ty, 5.0, 0.35, lbl, size=12, color=color, bold=True)


msg(gm_x, 2.4, fol_x, 2.9, "Sync   (t1 on TX, t2 on RX)", 4.0, 2.35)
msg(gm_x, 3.2, fol_x, 3.4, "Follow_Up  (carries t1)",     4.0, 3.05)
msg(fol_x, 3.9, gm_x, 4.3, "Delay_Req   (t3 on TX, t4 on RX)", 4.0, 4.05)
msg(gm_x, 4.7, fol_x, 5.0, "Delay_Resp  (carries t4)",    4.0, 4.55)

for x, y, lab in [(gm_x, 2.4, "t1"), (fol_x, 2.9, "t2"),
                  (fol_x, 3.9, "t3"), (gm_x, 4.3, "t4")]:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(x-0.08), Inches(y-0.08), Inches(0.16), Inches(0.16))
    d.fill.solid(); d.fill.fore_color.rgb = WARN; d.line.color.rgb = WARN
    label(s, x+0.15, y-0.18, 0.6, 0.3, lab, size=12, bold=True, color=WARN)

box(s, 0.6, 6.0, 12.0, 0.55,
    "mean_path_delay = ((t2 − t1) + (t4 − t3)) / 2       "
    "offset_from_master = (t2 − t1) − mean_path_delay",
    fill=LIGHT, fc=NAVY, size=13)

# ============================================================================
# 5. PTP Theory — Hardware timestamping
# ============================================================================
add_slide(L_CONTENT, "PTP — Hardware Timestamping", [
    "Where the timestamp is taken dictates the achievable accuracy",
    "Software timestamping (at socket layer):",
    (1, "Jitter from kernel scheduling, stack queues, cache effects"),
    (1, "Typical accuracy: tens of microseconds to low milliseconds"),
    "Hardware timestamping (at PHY/MAC Start-of-Frame):",
    (1, "Timestamp latched on the SFD bit — ideal reference point"),
    (1, "Removes OS/driver jitter from the measurement"),
    (1, "Typical accuracy: tens of nanoseconds, limited by counter resolution"),
    "IEEE 1588 assumes path symmetry — PTP cannot measure asymmetry directly",
    (1, "Any TX/RX pipeline asymmetry becomes a constant offset → must be calibrated"),
])

# ============================================================================
# 6. PTP Theory — Error sources
# ============================================================================
add_slide(L_CONTENT, "PTP — Error Sources", [
    "Oscillator rate offset and wander (short- and long-term)",
    "Quantisation of hardware timestamps (timestamp counter tick period)",
    "Asymmetric TX vs RX datapath latency — not observable by PTP math",
    "Link-layer residence time variance (queuing, media access)",
    (1, "Especially relevant on shared media such as 10BASE-T1S PLCA"),
    "Temperature drift of crystals → rate changes slowly over minutes/hours",
    "Servo dynamics:",
    (1, "Fast servo → tracks wander, amplifies measurement noise"),
    (1, "Slow servo → quiet but sluggish against temperature transients"),
])

# ============================================================================
# 7. 10BASE-T1S — PHY overview
# ============================================================================
add_slide(L_CONTENT, "10BASE-T1S — PHY Overview", [
    "10 Mbit/s full-duplex-equivalent over a single unshielded twisted pair",
    "Two physical topologies supported:",
    (1, "Point-to-point — two nodes sharing one pair"),
    (1, "Multidrop — many nodes on one shared segment (up to 8+ specified)"),
    "Targets automotive in-vehicle networking and industrial sensor aggregation",
    "Cable reach up to ~25 m on a shared segment, low EMI, low cost wiring",
    "Media access arbitration via PLCA (Physical Layer Collision Avoidance):",
    (1, "Time-slotted, deterministic round-robin transmit opportunities"),
    (1, "No CSMA/CD backoff → bounded worst-case access latency"),
])

# ============================================================================
# 8. 10BASE-T1S — PTP implications
# ============================================================================
add_slide(L_CONTENT, "10BASE-T1S — PTP Implications", [
    "PTP works over 10BASE-T1S the same way as over ordinary Ethernet",
    "Favourable properties:",
    (1, "PLCA bounds media-access jitter — residence time is predictable"),
    (1, "Physical medium is symmetric by construction (same pair, same direction)"),
    "Challenges:",
    (1, "10 Mbit/s → one byte takes 800 ns, lower timestamp resolution than 1 GbE"),
    (1, "PLCA slot wait adds variable queuing before TX"),
    (2, "→ hardware SFD timestamping removes this from the PTP measurement"),
    (1, "MAC-PHY vendor implementations vary — must characterise the specific chip"),
])

# ============================================================================
# 9. LAN8651 — Device overview
# ============================================================================
add_slide(L_CONTENT, "LAN8651 — Device Overview", [
    "Microchip 10BASE-T1S MAC-PHY with SPI host interface",
    "Integrated MAC + PHY in a single package",
    (1, "MCU needs no on-chip Ethernet MAC"),
    (1, "OPEN Alliance TC6 SPI protocol — data chunks over SPI"),
    "Supports PLCA for multidrop topologies",
    "Hardware timestamp engine latches TX-SFD and RX-SFD events",
    (1, "Raw material for PTP — accuracy depends on SFD latch point"),
    "Interrupt-driven host signalling via nIRQ",
    "Pin-count and BOM optimised for automotive and industrial edge nodes",
])

# ============================================================================
# 10. LAN8651 — Timestamping signal path
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "LAN8651 — Timestamping Signal Path")

# flow blocks
box(s, 0.6, 2.1, 2.1, 1.0, "T1S wire\nSFD bit", fill=OK, size=13)
box(s, 3.0, 2.1, 2.1, 1.0, "LAN8651 PHY\nSFD detect", fill=ACCENT, size=13)
box(s, 5.4, 2.1, 2.1, 1.0, "LAN8651 MAC\ntimestamp latch", fill=ACCENT, size=13)
box(s, 7.8, 2.1, 2.1, 1.0, "SPI readout\nto MCU", fill=NAVY, size=13)
box(s, 10.2, 2.1, 2.1, 1.0, "nIRQ\nto MCU", fill=WARN, size=13)
for x in [2.7, 5.1, 7.5, 9.9]:
    line(s, x, 2.6, x+0.3, 2.6, color=GREY, width=2, arrow=True)

# annotations
notes = [
    "• TX path: MCU writes frame over SPI → MAC transmits → SFD on wire = t1 latched",
    "• RX path: SFD on wire = t2 latched → MAC stores timestamp → nIRQ → MCU reads over SPI",
    "• TX: ~800 µs from SFD-on-wire to the TX-timestamp-available nIRQ  (pipeline + SPI)",
    "• RX: ~10 ms from SFD-on-wire to the RX-nIRQ   (LAN8651 behaviour)",
    "• Both biases must be compensated outside PTP — protocol cannot see them",
]
tb = s.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12.0), Inches(3.5))
tf = tb.text_frame; tf.word_wrap = True
for i, n in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = n
    r.font.size = Pt(16); r.font.color.rgb = NAVY

# ============================================================================
# 11. LAN8651 — Calibrating the biases
# ============================================================================
add_slide(L_CONTENT, "LAN8651 — Calibrating the Biases", [
    "Project design: ISR-captured anchor pair  (wall-clock_ns, TC0_tick)",
    "Anchor offsets shift the wall-clock timestamp forward to the ISR moment:",
    (1, "PTP_GM_ANCHOR_OFFSET_NS  = 800 000 ns     (TX pipeline)"),
    (1, "PTP_FOL_ANCHOR_OFFSET_NS = 10 000 000 ns  (RX pipeline)"),
    "Calibration is empirical:",
    (1, "Measure cross-board edge delta D with cyclic_fire_hw_test.py"),
    (1, "new_offset = old_offset + D    (positive D means FOL lags GM)"),
    "Effect of the 10 ms compensation (commit deb2773):",
    (1, "Cross-board median delta:  +10 005 µs  →  +10 µs"),
])

# ============================================================================
# 12. Project Integration — firmware map
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Project Integration — Firmware Map")

box(s, 0.6, 1.7, 3.8, 1.15, "ptp_gm_task.c\nGrandmaster state machine", fill=NAVY, size=13)
box(s, 4.7, 1.7, 3.8, 1.15, "ptp_fol_task.c\nFollower PI servo",          fill=NAVY, size=13)
box(s, 8.6, 1.7, 3.8, 1.15, "ptp_clock.c\nNanosecond software clock",     fill=NAVY, size=13)
box(s, 0.6, 3.0, 11.8, 0.95,
    "drv_lan865x_api.c — SPI MAC-PHY driver + TX/RX timestamp capture",
    fill=ACCENT, size=13)
box(s, 0.6, 4.1, 11.8, 0.9,
    "Harmony 3 TCP/IP stack  (net_10base_t1s v1.4.3)",
    fill=GREY, size=12)
box(s, 0.6, 5.1, 11.8, 0.9,
    "ATSAME54P20A — TC0 @ 60 MHz timestamp counter, SPI master",
    fill=GREY, size=12)

label(s, 0.6, 6.1, 12.0, 0.3,
      "Grandmaster and Follower run identical firmware — role is configuration, not code.",
      size=14, italic=True, color=GREY)

# ============================================================================
# 13. Cyclic Fire — What is it?
# ============================================================================
add_slide(L_CONTENT, "Cyclic Fire — What Is It?", [
    "A firmware module that fires a periodic GPIO pulse train on PD10",
    (1, "Pin sits on the EXT1 Xplained-Pro header — directly scope-clippable"),
    "Each pulse edge is scheduled against the PTP software clock, not a local timer",
    (1, "Implemented with tfuture — single-shot timer re-armed from its own callback"),
    "The scheduling time is given in absolute PTP wall-clock nanoseconds",
    (1, "cyclic_fire_start(period_us, phase_anchor_ns)"),
    (1, "phase_anchor_ns = the PTP moment at which the very first edge should fire"),
    "Identical firmware, identical API call on both boards",
    (1, "→ a shared anchor_ns produces wire-level coincident edges across boards"),
    "A CPU-cheap, direct-from-wire way to observe whether two PTP clocks really agree",
])

# ============================================================================
# 14. Cyclic Fire — What can you achieve with it?
# ============================================================================
add_slide(L_CONTENT, "Cyclic Fire — What You Can Achieve", [
    "Sync quality measurement: cross-board rising-edge delta quantifies PTP accuracy",
    (1, "Median → static bias (calibration target)"),
    (1, "MAD / extrema → short-term wander, filter behaviour"),
    "Calibration feedback loop: run test → read D → update anchor offset → rerun",
    "PTP regression test: a single gated CI metric (|median| + MAD ≤ 100 µs)",
    "Sub-µs demonstration target: visible proof that two independent MCUs tick together",
    "Distributed actuation reference: template for real coordinated outputs",
    (1, "Sensor sampling on N nodes at the same nanosecond"),
    (1, "Time-slotted TX windows (TDMA over T1S)"),
    (1, "Synchronised ADC trigger / actuator pulse"),
    "Failure detection: PLCA starvation, MAC wedge, clock drift all show up as edge drift",
])

# ============================================================================
# 15. Cyclic Fire — API & Patterns
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Cyclic Fire — API & Output Patterns")

# API blocks
box(s, 0.5, 1.55, 12.3, 0.55,
    "bool cyclic_fire_start(uint32_t period_us, uint64_t phase_anchor_ns);",
    fill=NAVY, size=14)
box(s, 0.5, 2.15, 12.3, 0.55,
    "bool cyclic_fire_start_ex(uint32_t period_us, uint64_t phase_anchor_ns, cyclic_fire_pattern_t pattern);",
    fill=NAVY, size=13)
box(s, 0.5, 2.75, 6.0, 0.55, "void cyclic_fire_stop(void);",             fill=GREY, size=13)
box(s, 6.8, 2.75, 6.0, 0.55, "bool cyclic_fire_is_running(void);",       fill=GREY, size=13)
box(s, 0.5, 3.35, 6.0, 0.55, "uint64_t cyclic_fire_get_cycle_count(void);",  fill=GREY, size=13)
box(s, 6.8, 3.35, 6.0, 0.55, "uint64_t cyclic_fire_get_missed_count(void);", fill=GREY, size=13)

# Pattern description
box(s, 0.5, 4.15, 6.0, 0.45, "PATTERN_SQUARE  — 50/50 rectangle (default)", fill=ACCENT, size=13)
box(s, 6.8, 4.15, 6.0, 0.45, "PATTERN_MARKER  — 1-high + 4-low isolated pulse", fill=ACCENT, size=13)

# Notes
notes = [
    "• period_us is the FULL rectangle period — callback fires every period_us/2",
    "• Default 1000 µs → 1 kHz square wave",
    "• Periods below ~400 µs not recommended (half-period ≈ tfuture spin threshold)",
    "• MARKER pattern isolates the rising edge — makes 'which board fires first?' unambiguous",
    "• CLI wrappers: cyclic_start, cyclic_start_marker, cyclic_start_free, cyclic_stop, cyclic_status",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
for i, n in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = n
    r.font.size = Pt(14); r.font.color.rgb = NAVY

# ============================================================================
# Cyclic Fire — Wire-Level View (Picture2)
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Cyclic Fire — Wire-Level View")
# image: 4.0 in tall, aspect ≈ 1.66 → 6.65 in wide, centred horizontally
s.shapes.add_picture(r"C:\work\ptp\check\net_10base_t1s\Picture2.png",
                     Inches(3.34), Inches(1.45), height=Inches(4.0))
# caption / explanation below the image
tb = s.shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
cap_lines = [
    ("Saleae Logic 2 capture — Ch0 (top) = Grandmaster PD10,  Ch1 (bottom) = Follower PD10.", True),
    ("MARKER pattern: 1 high half-period + 9 low half-periods → one isolated rising edge every 5 full periods.", False),
    ("Both channels fire coincident edges over 0.3 s — visible evidence that the two boards share a single PTP time base.", False),
]
for i, (text, bold) in enumerate(cap_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.color.rgb = NAVY; r.font.bold = bold

# ============================================================================
# Cyclic Fire — Zoom-In: Cross-Board Edge Delta (Picture3)
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Cyclic Fire — Zoom-In: Cross-Board Edge Delta")
s.shapes.add_picture(r"C:\work\ptp\check\net_10base_t1s\Picture3.png",
                     Inches(3.34), Inches(1.45), height=Inches(4.0))
tb = s.shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
cap_lines = [
    ("Zoomed Logic 2 view of a single MARKER pulse pair  (period_us = 5000 → 20 % duty, 40 Hz).", True),
    ("Cursor tool ΔT = 34.98 µs — time between the Grandmaster rising edge (Ch0) and the Follower rising edge (Ch1).", False),
    ("Well inside the 100 µs product window — this single wire-level frame corroborates the automated 10 µs median statistic.", False),
]
for i, (text, bold) in enumerate(cap_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = text
    r.font.size = Pt(14); r.font.color.rgb = NAVY; r.font.bold = bold

# ============================================================================
# PTP_CLOCK — The Application API for PTP Time
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "PTP_CLOCK — Application API for PTP Time")

box(s, 0.5, 1.55, 12.3, 0.55,
    "uint64_t PTP_CLOCK_GetTime_ns(void);         // current wall-clock, interpolated",
    fill=NAVY, size=13)
box(s, 0.5, 2.15, 12.3, 0.55,
    "bool     PTP_CLOCK_IsValid(void);            // true once first anchor was set",
    fill=NAVY, size=13)
box(s, 0.5, 2.75, 12.3, 0.55,
    "int32_t  PTP_CLOCK_GetDriftPPB(void);        // measured MCU vs PTP rate offset",
    fill=GREY, size=13)
box(s, 0.5, 3.35, 12.3, 0.55,
    "void     PTP_CLOCK_ForceSet(uint64_t ns);    // standalone set (no PTP sync)",
    fill=GREY, size=13)

notes = [
    "• Any context — no SPI, no mutex, no blocking.  Safe from ISR and main loop.",
    "• Interpolates between anchors via TC0 (60 MHz) + drift IIR filter.",
    "• One anchor is refreshed on every PTP Sync (≈ 8 Hz).",
    "• Returns 0 until the first anchor has been set — always check PTP_CLOCK_IsValid() first.",
    "• GM and Follower expose the same API — application code is role-agnostic.",
    "• Resolution ≈ 16.7 ns (TC0 tick); short-term jitter dominated by filter wander.",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(4.15), Inches(12.3), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
for i, n in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = n
    r.font.size = Pt(14); r.font.color.rgb = NAVY

# ============================================================================
# 17. PTP_CLOCK — Event Logging Pattern
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "PTP_CLOCK — Logging Firmware Events")

# Minimal, readable code snippet
code_lines = [
    "#include \"ptp_clock.h\"",
    "",
    "static void on_frame_received(void) {",
    "    uint64_t t_ns = PTP_CLOCK_GetTime_ns();     // synchronised wall-clock",
    "    SYS_CONSOLE_PRINT(\"RX %llu.%09llu  seq=%u\\r\\n\",",
    "                      t_ns / 1000000000ull,",
    "                      t_ns % 1000000000ull,",
    "                      rx_seq);",
    "}",
    "",
    "void actuator_fire(uint64_t at_ns) {",
    "    while (PTP_CLOCK_GetTime_ns() < at_ns) { /* tight spin */ }",
    "    GPIO_PinSet(ACT_PIN);",
    "}",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.8))
tf = tb.text_frame; tf.word_wrap = True
for i, ln in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = ln if ln else " "
    r.font.name = "Consolas"
    r.font.size = Pt(14)
    r.font.color.rgb = NAVY

notes = [
    "• Pattern: call PTP_CLOCK_GetTime_ns() at the moment of interest, store alongside the event.",
    "• Log to console, ring buffer, or CLI — the timestamp is portable across boards.",
    "• For event correlation across nodes: compare two logs directly; no relative-time conversion.",
    "• For scheduled actions: give the target time to the action, spin or sleep until reached.",
    "• Always guard with PTP_CLOCK_IsValid() during startup to avoid 0-timestamps.",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
for i, n in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = n
    r.font.size = Pt(13); r.font.color.rgb = NAVY

# ============================================================================
# 18. PTP_CLOCK — Cross-Board Correlation
# ============================================================================
add_slide(L_CONTENT, "PTP_CLOCK — Cross-Board Correlation", [
    "Every event stamped with PTP_CLOCK_GetTime_ns() lives on a common timeline",
    "Typical workflows the API enables:",
    (1, "Distributed trace — collect logs from N nodes, sort by PTP timestamp"),
    (1, "Race analysis — precise order of events on nodes A, B, C"),
    (1, "Latency measurement — PTP_CLOCK_ns at TX on A vs RX on B"),
    (1, "Scheduled output — compute target ns, hand to application, pin fires synchronously"),
    (1, "Drift surveillance — log PTP_CLOCK_GetDriftPPB() over long runs"),
    "Accuracy envelope (current firmware):",
    (1, "Static bias between boards: ~10 µs after calibration"),
    (1, "Short-term wander: ~40 µs RMS within a 0.7 s window"),
    (1, "Long-term rate agreement: 1.2 ppm over 60 s"),
    "Cyclic Fire is the canonical consumer of this API — any feature using PTP time follows the same pattern",
])

# ============================================================================
# 14. Results — methodology
# ============================================================================
add_slide(L_CONTENT, "Results — Methodology", [
    "Saleae Logic 8 — Ch0 = GM board PD10, Ch1 = Follower board PD10",
    "cyclic_fire_hw_test.py orchestrates the full run:",
    (1, "Reset both boards → wait for FINE PTP state"),
    (1, "Issue cyclic_start with a shared future anchor_ns"),
    (1, "Capture ~0.7 s via the Logic 2 scripting socket"),
    (1, "Extract rising edges, pair cross-board edges, compute delta distribution"),
    "Reported statistics: median, MAD, extrema, per-edge CSV dump",
    "Pass gate:  |median| ≤ 50 µs  AND  MAD ≤ 50 µs    →    sum ≤ 100 µs window",
    "Spot-check: Logic 2 cursor measurements cross-validated the script output",
])

# ============================================================================
# 15. Results — numbers
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Results — Measurement Summary")

box(s, 0.7, 1.9, 3.7, 1.8, "Median\n+10 µs",  fill=OK,     size=26)
box(s, 4.7, 1.9, 3.7, 1.8, "MAD\n38 µs",      fill=ACCENT, size=26)
box(s, 8.7, 1.9, 3.7, 1.8, "Verdict\nPASS",   fill=OK,     size=26)

notes = [
    "• Before RX-IRQ compensation:  median +10 005 µs  →  FAIL",
    "• After the fix:  median +10 µs  →  inside the 100 µs product window",
    "• Long-term (60 s) cross-board rate residual:  1.2 ppm",
    "• MAD dominated by short-term drift-filter wander, not stack asymmetry",
    "• Drift IIR raised from N=32 to N=128 → halved filter stddev",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12.0), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
for i, n in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = n
    r.font.size = Pt(17); r.font.color.rgb = NAVY

# ============================================================================
# Results — Cursor Measurement (Picture3)
# ============================================================================
# ============================================================================
# Open issues & outlook
# ============================================================================
# ============================================================================
# Part 7 — Reproducing the project: Prerequisites & Clone
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Reproducing the Project — Prerequisites")

# Hardware block
box(s, 0.5, 1.5, 12.3, 0.55, "Hardware", fill=NAVY, size=14)
label(s, 0.7, 2.1, 12.0, 0.6,
      "2 × ATSAME54 Curiosity Ultra  +  2 × LAN865x click board  —  connected pair-to-pair over 10BASE-T1S",
      size=14, color=NAVY)
label(s, 0.7, 2.5, 12.0, 0.4,
      "USB (EDBG) cables from each board to the PC for flashing and for the serial console (115200 8N1).",
      size=13, color=GREY)

# Software block
box(s, 0.5, 3.15, 12.3, 0.55, "Software on the PC", fill=NAVY, size=14)
bullets_notes = [
    "MPLAB XC32 v4.60 or v5.x  (default path  C:\\Program Files\\Microchip\\xc32\\)",
    "MPLAB X IDE / MDB  (used by flash.py to program the boards)",
    "CMake ≥ 4.1  +  Ninja  (both on PATH)",
    "Python 3.9+   with   pip install pyserial   (plus python-pptx only if you rebuild this deck)",
    "Two terminal windows — one per board — e.g. PuTTY or Tera Term  (115200 8N1)",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(12.0), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
for i, line_txt in enumerate(bullets_notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = "• " + line_txt
    r.font.size = Pt(14); r.font.color.rgb = NAVY

# Clone block
box(s, 0.5, 6.05, 12.3, 0.55, "Clone", fill=ACCENT, size=13)
tb = s.shapes.add_textbox(Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
cmds = [
    "git clone https://github.com/zabooh/net_10base_t1s.git",
    "cd net_10base_t1s\\apps\\tcpip_iperf_lan865x\\firmware\\tcpip_iperf_lan865x.X",
]
for i, c in enumerate(cmds):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = c
    r.font.name = "Consolas"; r.font.size = Pt(13); r.font.color.rgb = NAVY

# ============================================================================
# Part 7 — Reproducing the project: Setup, Build, Flash, Run
# ============================================================================
s = prs.slides.add_slide(L_TITLE_ONLY)
set_title(s, "Reproducing the Project — Setup, Build, Run")

# Quick-path banner — skip the build entirely, just flash the checked-in HEX
box(s, 0.5, 1.4, 12.3, 0.55,
    "Quick path — no build required (just run the demo)", fill=OK, size=14)
quick_lines = [
    "After  git clone  and  python setup_flasher.py,  simply run:",
    "    python flash.py",
    "flash.py finds the newest pre-built HEX under  out/…/image/*.hex  (checked into the repo) "
    "and programmes it onto both boards — no XC32 / CMake / build step needed.",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
for i, lt in enumerate(quick_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = lt
    r.font.size = Pt(13); r.font.color.rgb = NAVY
    if lt.lstrip() != lt:
        r.font.name = "Consolas"

# Step 1: one-time setup
box(s, 0.5, 3.35, 12.3, 0.45,
    "Full rebuild  —  1.  One-time tool setup  (run once per machine)", fill=NAVY, size=13)
tb = s.shapes.add_textbox(Inches(0.7), Inches(3.85), Inches(12.0), Inches(1.1))
tf = tb.text_frame; tf.word_wrap = True
cmds1 = [
    "python setup_compiler.py      # pick the installed XC32 version",
    "python setup_flasher.py       # assign Board 1 (GM) and Board 2 (FOL) to their EDBG debuggers",
    "python setup_debug.py         # fix SAME54_DFP tool-pack (needed for VS Code debugging)",
]
for i, c in enumerate(cmds1):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = c
    r.font.name = "Consolas"; r.font.size = Pt(12); r.font.color.rgb = NAVY

# Step 2: build and flash
box(s, 0.5, 4.95, 12.3, 0.45, "2.  Build and flash both boards", fill=NAVY, size=13)
tb = s.shapes.add_textbox(Inches(0.7), Inches(5.45), Inches(12.0), Inches(0.85))
tf = tb.text_frame; tf.word_wrap = True
cmds2 = [
    "build.bat                     # incremental build   (build.bat rebuild for a clean build)",
    "python flash.py               # flash Board 1 and Board 2 in sequence",
]
for i, c in enumerate(cmds2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = c
    r.font.name = "Consolas"; r.font.size = Pt(12); r.font.color.rgb = NAVY

# Step 3: run the demo
box(s, 0.5, 6.30, 12.3, 0.45, "3.  Run the Cyclic Fire demo", fill=ACCENT, size=13)
tb = s.shapes.add_textbox(Inches(0.7), Inches(6.80), Inches(12.0), Inches(1.1))
tf = tb.text_frame; tf.word_wrap = True
demo_lines = [
    "Open two terminals (115200 8N1).  Wait for PTP state = FINE.  On each:",
    "    cyclic_start_marker 5000 <anchor_ns>    # or:  python cyclic_fire_hw_test.py --marker --period-us 5000",
]
for i, lt in enumerate(demo_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = lt
    r.font.size = Pt(12); r.font.color.rgb = NAVY
    if lt.lstrip() != lt:
        r.font.name = "Consolas"

# ============================================================================
# Open Issues & Outlook
# ============================================================================
add_slide(L_CONTENT, "Open Issues & Outlook", [
    "Short-term wander within 0.7 s capture windows: 50…200 µs/s drift bursts",
    (1, "Drift filter random-walks (lag-1 autocorrelation > 0.9)"),
    "Paths forward:",
    (1, "Longer averaging window (N = 256) — trades response speed for variance"),
    (1, "Median-of-N filter instead of IIR for outlier resilience"),
    (1, "Hardware-scheduled GPIO anchored to the LAN8651 1PPS / MAC timer"),
    (2, "→ bypasses software PTP_CLOCK, opens the sub-microsecond regime"),
    "Known failure mode: LAN8651 MAC occasionally wedges after long runs",
    (1, "Soft reset insufficient; USB power-cycle required"),
    "References in-tree: README_PTP §§4.3, 11, 12  and  README_NTP §8",
])

prs.save(PATH)
print(f"OK — wrote {len(prs.slides)} slides to {PATH}")
